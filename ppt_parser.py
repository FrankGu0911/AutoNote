import os
import pptx
from pptx import Presentation
from typing import Dict, List, Any, Optional
import base64
from io import BytesIO

class PPTParser:
    """
    PPT解析类，负责从PPT文件中提取内容
    """
    def __init__(self):
        self.presentation = None
        self.slides_data = []
    
    def load_presentation(self, file_path: str) -> bool:
        """
        加载PPT文件
        
        Args:
            file_path: PPT文件路径
            
        Returns:
            bool: 是否成功加载
        """
        if not os.path.exists(file_path):
            print(f"文件不存在: {file_path}")
            return False
        
        try:
            self.presentation = Presentation(file_path)
            return True
        except Exception as e:
            print(f"加载PPT文件失败: {e}")
            return False
    
    def extract_text_from_shape(self, shape) -> str:
        """从形状中提取文本"""
        text = ""
        if hasattr(shape, "text"):
            text = shape.text
        elif shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text += run.text
        return text
    
    def extract_image_data(self, shape) -> Optional[str]:
        """提取图片并转换为base64字符串"""
        if not hasattr(shape, "image"):
            return None
        
        try:
            image_bytes = shape.image.blob
            image_type = shape.image.content_type
            encoded = base64.b64encode(image_bytes).decode('utf-8')
            return f"data:{image_type};base64,{encoded}"
        except Exception as e:
            print(f"提取图片数据失败: {e}")
            return None
    
    def parse_slide(self, slide, slide_index: int) -> Dict[str, Any]:
        """
        解析单个幻灯片
        
        Args:
            slide: 幻灯片对象
            slide_index: 幻灯片索引
            
        Returns:
            Dict: 包含幻灯片解析数据的字典
        """
        slide_data = {
            "index": slide_index,
            "title": "",
            "content": [],
            "images": []
        }
        
        # 提取标题
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text
        
        # 提取所有形状中的内容
        for shape in slide.shapes:
            # 跳过标题形状，因为已经处理过
            if shape == slide.shapes.title:
                continue
            
            # 提取文本
            text = self.extract_text_from_shape(shape)
            if text:
                slide_data["content"].append(text)
            
            # 提取图片
            image_data = self.extract_image_data(shape)
            if image_data:
                slide_data["images"].append({
                    "data": image_data,
                    "alt_text": shape.alt_text if hasattr(shape, "alt_text") else ""
                })
        
        return slide_data
    
    def parse_presentation(self) -> List[Dict[str, Any]]:
        """
        解析整个PPT演示文稿
        
        Returns:
            List[Dict]: 包含所有幻灯片数据的列表
        """
        if not self.presentation:
            print("未加载PPT文件")
            return []
        
        self.slides_data = []
        
        for i, slide in enumerate(self.presentation.slides):
            slide_data = self.parse_slide(slide, i)
            self.slides_data.append(slide_data)
        
        return self.slides_data
    
    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        """
        处理PPT文件并返回解析后的数据
        
        Args:
            file_path: PPT文件路径
            
        Returns:
            List[Dict]: 包含所有幻灯片数据的列表
        """
        if self.load_presentation(file_path):
            return self.parse_presentation()
        return []


if __name__ == "__main__":
    # 测试代码
    parser = PPTParser()
    slides = parser.process_file("test.pptx")
    print(f"解析了 {len(slides)} 张幻灯片")
    for slide in slides:
        print(f"幻灯片 {slide['index']+1}: {slide['title']}")
        print(f"内容: {slide['content']}")
        print(f"图片数量: {len(slide['images'])}")
        print("-" * 50) 