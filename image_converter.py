import os
import fitz  # PyMuPDF
from PIL import Image
import io
import tempfile
import shutil
from pptx import Presentation
from typing import List, Dict, Any, Optional, Tuple
import config
from utils.helpers import ensure_directory_exists, is_ppt_file, is_pdf_file

class ImageConverter:
    """
    文档图像转换器，负责将PPT或PDF文件转换为高质量图像
    """
    def __init__(self, temp_dir: str = None, images_dir: str = None):
        """
        初始化图像转换器
        
        Args:
            temp_dir: 临时图像存储目录
            images_dir: 笔记图像存储目录
        """
        self.temp_dir = temp_dir or config.IMAGE_TEMP_DIR
        self.images_dir = images_dir or os.path.join(config.DEFAULT_OUTPUT_DIR, "images")
        ensure_directory_exists(self.temp_dir)
        ensure_directory_exists(self.images_dir)
        self.image_paths = []
        self.markdown_image_paths = []
        self.dpi = config.IMAGE_DPI
        self.image_format = config.IMAGE_FORMAT
        self.image_quality = config.IMAGE_QUALITY
    
    def convert_pdf_to_images(self, file_path: str, save_for_markdown: bool = False) -> List[Dict[str, Any]]:
        """
        将PDF文件转换为图像
        
        Args:
            file_path: PDF文件路径
            save_for_markdown: 是否保存图像用于Markdown引用
            
        Returns:
            List[Dict]: 包含图像路径和元数据的列表
        """
        self.image_paths = []
        self.markdown_image_paths = []
        result = []
        
        try:
            # 创建文档特定的图像目录
            if save_for_markdown:
                base_name = os.path.basename(file_path)
                doc_name = os.path.splitext(base_name)[0]
                doc_images_dir = os.path.join(self.images_dir, doc_name)
                ensure_directory_exists(doc_images_dir)
            
            # 打开PDF文件
            pdf_document = fitz.open(file_path)
            
            # 遍历每一页
            for page_index in range(len(pdf_document)):
                page = pdf_document[page_index]
                
                # 渲染页面为图像
                pix = page.get_pixmap(dpi=self.dpi)
                
                # 创建图像路径
                image_path = os.path.join(self.temp_dir, f"page_{page_index+1}.{self.image_format}")
                
                # 保存图像
                if self.image_format.lower() == "png":
                    pix.save(image_path)
                else:
                    # 转换为PIL图像以支持其他格式
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img.save(image_path, quality=self.image_quality)
                
                self.image_paths.append(image_path)
                
                # 提取基本文本作为辅助信息
                text = page.get_text()
                
                # 尝试提取标题
                title = self._extract_page_title(text, page_index)
                
                # 提取页面中的图片并保存
                page_images = []
                if save_for_markdown:
                    page_images = self._extract_and_save_pdf_images(
                        pdf_document, page, page_index, doc_images_dir
                    )
                
                # 创建页面数据
                page_data = {
                    "index": page_index,
                    "title": title,
                    "image_path": image_path,
                    "text_hint": text[:500] if text else "",  # 仅保留部分文本作为辅助提示
                    "images": page_images
                }
                
                result.append(page_data)
            
            pdf_document.close()
            
        except Exception as e:
            print(f"PDF转图像失败: {e}")
        
        return result
    
    def _extract_and_save_pdf_images(self, pdf_document, page, page_index: int, output_dir: str) -> List[Dict[str, Any]]:
        """
        提取PDF页面中的图片并保存用于Markdown引用
        
        Args:
            pdf_document: PDF文档对象
            page: PDF页面对象
            page_index: 页面索引
            output_dir: 输出目录
            
        Returns:
            List[Dict]: 包含图片信息的列表
        """
        page_images = []
        
        try:
            # 获取页面上的图片
            image_list = page.get_images(full=True)
            
            for img_index, img_info in enumerate(image_list):
                xref = img_info[0]
                base_image = pdf_document.extract_image(xref)
                
                if base_image:
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # 创建Markdown图片路径和文件名
                    md_image_name = f"page{page_index+1}-img{img_index+1}.{image_ext}"
                    md_image_path = os.path.join(output_dir, md_image_name)
                    
                    # 保存图片
                    with open(md_image_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    # 相对路径，用于Markdown引用
                    rel_path = os.path.join("images", os.path.basename(output_dir), md_image_name)
                    
                    # 添加到图片列表
                    page_images.append({
                        "index": img_index,
                        "page": page_index,
                        "path": md_image_path,
                        "relative_path": rel_path.replace("\\", "/")
                    })
                    
                    self.markdown_image_paths.append(md_image_path)
        except Exception as e:
            print(f"提取PDF图片失败: {e}")
        
        return page_images
    
    def convert_ppt_to_images(self, file_path: str, save_for_markdown: bool = False) -> List[Dict[str, Any]]:
        """
        将PPT文件转换为图像
        
        Args:
            file_path: PPT文件路径
            save_for_markdown: 是否保存图像用于Markdown引用
            
        Returns:
            List[Dict]: 包含图像路径和元数据的列表
        """
        self.image_paths = []
        self.markdown_image_paths = []
        result = []
        
        try:
            # 创建文档特定的图像目录
            if save_for_markdown:
                base_name = os.path.basename(file_path)
                doc_name = os.path.splitext(base_name)[0]
                doc_images_dir = os.path.join(self.images_dir, doc_name)
                ensure_directory_exists(doc_images_dir)
            
            # 打开PPT文件
            if save_for_markdown:
                # 直接从PPT提取图片
                presentation = Presentation(file_path)
                slides_with_images = self._extract_and_save_ppt_images(presentation, doc_images_dir)
            
            # 使用PDF作为中间格式
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_pdf:
                tmp_pdf_path = tmp_pdf.name
            
            self._convert_ppt_to_pdf(file_path, tmp_pdf_path)
            
            # 然后使用现有的PDF转图像功能
            result = self.convert_pdf_to_images(tmp_pdf_path, False)  # 不从PDF中提取图片
            
            # 如果保存了PPT的图片，将其添加到结果中
            if save_for_markdown and slides_with_images:
                for i, slide_data in enumerate(result):
                    if i < len(slides_with_images):
                        slide_data["images"] = slides_with_images[i]
            
            # 删除临时PDF文件
            os.unlink(tmp_pdf_path)
            
            # 对于PPT，尝试从原始文件中提取标题
            self._update_titles_from_ppt(file_path, result)
            
        except Exception as e:
            print(f"PPT转图像失败: {e}")
        
        return result
    
    def _extract_and_save_ppt_images(self, presentation, output_dir: str) -> List[List[Dict[str, Any]]]:
        """
        提取PPT中的图片并保存用于Markdown引用
        
        Args:
            presentation: PPT演示文稿对象
            output_dir: 输出目录
            
        Returns:
            List[List[Dict]]: 包含幻灯片图片信息的嵌套列表
        """
        slides_with_images = []
        
        try:
            for slide_index, slide in enumerate(presentation.slides):
                slide_images = []
                img_index = 0
                
                for shape in slide.shapes:
                    # 检查是否是图片
                    if hasattr(shape, "image") and shape.image:
                        try:
                            # 创建Markdown图片路径和文件名
                            image_ext = "png"  # 默认扩展名
                            content_type = shape.image.content_type
                            if content_type:
                                if "jpeg" in content_type or "jpg" in content_type:
                                    image_ext = "jpg"
                                elif "png" in content_type:
                                    image_ext = "png"
                                elif "gif" in content_type:
                                    image_ext = "gif"
                            
                            md_image_name = f"slide{slide_index+1}-img{img_index+1}.{image_ext}"
                            md_image_path = os.path.join(output_dir, md_image_name)
                            
                            # 保存图片
                            with open(md_image_path, "wb") as img_file:
                                img_file.write(shape.image.blob)
                            
                            # 相对路径，用于Markdown引用
                            rel_path = os.path.join("images", os.path.basename(output_dir), md_image_name)
                            
                            # 添加到图片列表
                            slide_images.append({
                                "index": img_index,
                                "page": slide_index,
                                "path": md_image_path,
                                "relative_path": rel_path.replace("\\", "/")
                            })
                            
                            self.markdown_image_paths.append(md_image_path)
                            img_index += 1
                        except Exception as e:
                            print(f"提取幻灯片 {slide_index+1} 图片失败: {e}")
                
                slides_with_images.append(slide_images)
        except Exception as e:
            print(f"提取PPT图片失败: {e}")
        
        return slides_with_images
    
    def _convert_ppt_to_pdf(self, ppt_path: str, pdf_path: str) -> bool:
        """
        将PPT转换为PDF（使用外部工具LibreOffice或Microsoft PowerPoint或提示用户）
        
        Args:
            ppt_path: PPT文件路径
            pdf_path: 输出PDF路径
            
        Returns:
            bool: 转换是否成功
        """
        # 方法1: 尝试使用Microsoft PowerPoint (Windows)
        try:
            import comtypes.client
            import os
            
            # PowerPoint常量
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = True
            
            # 转换常量
            ppSaveAsPDF = 32  # 导出为PDF格式的常量
            
            # 获取绝对路径
            abs_ppt_path = os.path.abspath(ppt_path)
            abs_pdf_path = os.path.abspath(pdf_path)
            
            # 打开PPT文件
            presentation = powerpoint.Presentations.Open(abs_ppt_path)
            
            # 保存为PDF
            presentation.SaveAs(abs_pdf_path, ppSaveAsPDF)
            
            # 关闭文件和应用程序
            presentation.Close()
            powerpoint.Quit()
            
            if os.path.exists(abs_pdf_path):
                return True
                
        except Exception as e:
            print(f"使用PowerPoint转换失败: {e}")
        
        # 方法2: 尝试使用LibreOffice转换（如果安装了的话）
        try:
            import subprocess
            result = subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", 
                 os.path.dirname(pdf_path), ppt_path],
                capture_output=True,
                text=True
            )
            
            if result.returncode == 0:
                # 重命名输出文件
                base_name = os.path.basename(ppt_path)
                file_name = os.path.splitext(base_name)[0]
                orig_pdf = os.path.join(os.path.dirname(pdf_path), f"{file_name}.pdf")
                os.rename(orig_pdf, pdf_path)
                return True
            
            raise Exception("LibreOffice转换失败")
            
        except Exception as e:
            print(f"无法自动转换PPT到PDF: {e}")
            print("请手动将PPT转换为PDF后重试，或使用文本模式处理。")
            return False
    
    def _update_titles_from_ppt(self, ppt_path: str, pages_data: List[Dict[str, Any]]) -> None:
        """
        从原始PPT中提取标题并更新页面数据
        
        Args:
            ppt_path: PPT文件路径
            pages_data: 页面数据列表
        """
        try:
            presentation = Presentation(ppt_path)
            
            for i, slide in enumerate(presentation.slides):
                if i < len(pages_data) and slide.shapes.title:
                    pages_data[i]["title"] = slide.shapes.title.text
        except:
            pass
    
    def _extract_page_title(self, text: str, page_index: int) -> str:
        """
        尝试从页面文本中提取标题
        
        Args:
            text: 页面文本
            page_index: 页面索引
            
        Returns:
            str: 提取的标题
        """
        if not text:
            return f"页面 {page_index + 1}"
        
        # 简单实现：尝试使用第一行作为标题
        lines = text.strip().split('\n')
        if lines:
            title = lines[0].strip()
            # 标题不应该太长
            if len(title) <= 100:
                return title
        
        return f"页面 {page_index + 1}"
    
    def convert_file_to_images(self, file_path: str, save_for_markdown: bool = False) -> List[Dict[str, Any]]:
        """
        根据文件类型将文件转换为图像
        
        Args:
            file_path: 文件路径
            save_for_markdown: 是否保存图像用于Markdown引用
            
        Returns:
            List[Dict]: 包含图像路径和元数据的列表
        """
        if is_pdf_file(file_path):
            return self.convert_pdf_to_images(file_path, save_for_markdown)
        elif is_ppt_file(file_path):
            return self.convert_ppt_to_images(file_path, save_for_markdown)
        else:
            raise ValueError(f"不支持的文件类型: {file_path}")
    
    def cleanup(self, keep_markdown_images: bool = True) -> None:
        """
        清理临时图像文件
        
        Args:
            keep_markdown_images: 是否保留用于Markdown的图片
        """
        # 删除临时分析图像
        for image_path in self.image_paths:
            if os.path.exists(image_path):
                try:
                    os.remove(image_path)
                except:
                    pass
        
        # 如果不保留Markdown图片，删除它们
        if not keep_markdown_images:
            for image_path in self.markdown_image_paths:
                if os.path.exists(image_path):
                    try:
                        os.remove(image_path)
                    except:
                        pass
    
    def get_image_placeholder_tag(self) -> str:
        """
        获取图片占位符标签
        
        Returns:
            str: 图片占位符标签
        """
        return "image:"


if __name__ == "__main__":
    # 测试代码
    converter = ImageConverter()
    import sys
    
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        print(f"正在转换文件: {file_path}")
        
        pages = converter.convert_file_to_images(file_path)
        print(f"已转换 {len(pages)} 页到图像")
        
        for page in pages:
            print(f"页面 {page['index']+1}: {page['title']}")
            print(f"图像路径: {page['image_path']}")
            print("-" * 50)
    else:
        print("请提供文件路径作为参数") 